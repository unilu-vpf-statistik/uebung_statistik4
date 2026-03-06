"""
Build S12: Anwendung II — Regression & Vorhersagen
Populates the scaffold PPTX with content from S12_Folieninhalte.md.
4 R-generated plots + text/code shapes.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *

BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_12_Regression.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_12_Regression_filled.pptx"
PLOTS = BASE / "local" / "Vorlesung" / "plots"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Sie wollen vorhersagen, wer gut',
              'abschneidet. Welche Information',
              'würden Sie nutzen?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Regression & Vorhersagen")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie..."],
        ["Regression als Verallge-", "meinerung des Gruppen-", "vergleichs verstehen"],
        ["Regressionskoeffizienten", "korrekt interpretieren", "(nicht nur Signifikanz)"],
        ["Interaktionseffekte", "konzeptuell erklären", "und visuell darstellen"],
    ]
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12, light=(i == 0))

    # --- Slide 8 (idx 7): Von ANOVA zu Regression ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)

    # Bridge diagram: ANOVA ← Linear Model → Regression
    add_hexagon(slide,
        left=Inches(0.5), top=Inches(2.0),
        width=HEX_LARGE_W, height=HEX_LARGE_H,
        text=["ANOVA:", "Kategorialer", "Prädiktor"],
        font_size=13, fill_scheme='accent3')
    add_arrow(slide,
        left=Inches(3.0), top=Inches(2.7),
        width=Inches(1.5), height=ARROW_H,
        direction='right')
    add_hexagon(slide,
        left=Inches(4.5), top=Inches(2.0),
        width=HEX_LARGE_W, height=HEX_LARGE_H,
        text=["Lineares Modell:", "y = b0 + b1*x + e"],
        font_size=13, bold=True)
    add_arrow(slide,
        left=Inches(7.0), top=Inches(2.7),
        width=Inches(1.5), height=ARROW_H,
        direction='right')
    add_hexagon(slide,
        left=Inches(8.5), top=Inches(2.0),
        width=HEX_LARGE_W, height=HEX_LARGE_H,
        text=["Regression:", "Metrischer", "Prädiktor"],
        font_size=13, fill_scheme='accent6')

    add_rounded_rect(slide,
        left=Inches(2.0), top=Inches(4.5),
        width=Inches(8.0), height=Inches(1.5),
        text=["ANOVA und Regression sind keine",
              "verschiedenen Verfahren -- sie sind",
              "Spezialfälle desselben Modells"],
        font_size=14)

    # --- Slide 9 (idx 8): Regressionsgleichung — R-Plot ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S12_F09_regression_linie.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 10 (idx 9): Koeffizienten interpretieren ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.0), height=HEX_CONTENT_H,
        text=["b1 = 0.8: Pro Lebensjahr steigt", "der WM-Score um 0.8 Punkte"],
        font_size=13)

    add_hexagon(slide,
        left=Inches(6.0), top=Inches(1.8),
        width=Inches(5.0), height=HEX_CONTENT_H,
        text=["b1 = -1.2: Pro Stunde mehr Stress", "sinkt der WM-Score um 1.2 Punkte"],
        font_size=13)

    # Common error
    add_hexagon(slide,
        left=Inches(0.5), top=Inches(3.5),
        width=Inches(5.0), height=HEX_CONTENT_H,
        text=["Falsch: Alter verursacht", "höheren WM-Score"],
        font_size=13, fill_scheme='accent2')
    add_hexagon(slide,
        left=Inches(6.0), top=Inches(3.5),
        width=Inches(5.0), height=HEX_CONTENT_H,
        text=["Richtig: Alter hängt mit", "höherem WM-Score zusammen"],
        font_size=13, fill_scheme='accent6')

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(5.5),
        width=Inches(9.0), height=Inches(1.2),
        text=["Regression zeigt Zusammenhang, nicht Kausalität!",
              "Grösse von b1 ist informativer als p"],
        font_size=13)

    # --- Slide 11 (idx 10): Modellgüte R² ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=["R-Quadrat = Anteil der Varianz in y,",
              "der durch das Modell erklärt wird"],
        font_size=14)

    examples = [
        ("R-Quadrat = 0", "Modell erklärt nichts"),
        ("R-Quadrat = .15", "15% erklärt (typisch Psychologie)"),
        ("R-Quadrat = 1", "Alles erklärt (unrealistisch)"),
    ]
    for j, (val, desc) in enumerate(examples):
        add_hexagon(slide,
            left=Inches(0.5 + j * 3.7), top=Inches(4.2),
            width=Inches(3.3), height=Inches(1.5),
            text=[val, desc], font_size=12, light=True)

    add_text_box(slide,
        left=Inches(2.0), top=Inches(6.0),
        width=Inches(8.0), height=Inches(0.8),
        text=["Einfache Regression: R-Quadrat = r-Quadrat (Korrelation zum Quadrat)"],
        font_size=13)

    # --- Slide 12 (idx 11): Residualdiagnostik — R-Plot ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S12_F12_diagnostik.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 13 (idx 12): Zusammenfassung Block A [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 14 (idx 13): Leitfrage Block B — R-Plot ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S12_F14_moderation.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 15 (idx 14): Multiple Regression ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(1.5),
        text=['lm(wm_score ~ alter + gruppe, data = daten)'],
        font_size=13)

    add_hexagon(slide,
        left=Inches(0.5), top=Inches(3.8),
        width=Inches(5.0), height=HEX_LARGE_H,
        text=["Koeffizienten sind adjustiert:",
              "Effekt von Alter kontrolliert",
              "für Gruppe (und umgekehrt)"],
        font_size=13)

    add_hexagon(slide,
        left=Inches(6.0), top=Inches(1.8),
        width=Inches(5.0), height=HEX_LARGE_H,
        text=["Einfach vs. Multipel:",
              "Koeffizient kann sich ändern",
              "→ Konfundierung erkennbar"],
        font_size=13, light=True)

    add_rounded_rect(slide,
        left=Inches(6.0), top=Inches(4.2),
        width=Inches(5.0), height=Inches(1.5),
        text=["Jeder Koeffizient zeigt den",
              "reinen Effekt, bereinigt",
              "um die anderen Prädiktoren"],
        font_size=13)

    # --- Slide 16 (idx 15): Interaktion — R-Plot ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S12_F16_interaktion.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 17 (idx 16): Interaktion in R ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(1.2),
        text=['modell_int <- lm(wm_score ~ alter * gruppe,',
              '                 data = daten)'],
        font_size=12)

    # Output explanation
    terms = [
        ("alter", "Effekt von Alter in KG (Referenz)"),
        ("gruppeEG", "Unterschied Achsenabschnitt EG vs KG"),
        ("alter:gruppeEG", "Unterschied Steigung = INTERAKTION"),
    ]
    for j, (term, desc) in enumerate(terms):
        scheme = 'accent2' if j == 2 else 'accent4'
        add_hexagon(slide,
            left=Inches(0.5), top=Inches(3.5 + j * 1.2),
            width=Inches(3.0), height=Inches(1.0),
            text=[term], font_size=12, fill_scheme=scheme)
        add_text_box(slide,
            left=Inches(4.0), top=Inches(3.5 + j * 1.2),
            width=Inches(4.0), height=Inches(1.0),
            text=[desc], font_size=13)

    add_code_box(slide,
        left=Inches(6.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(1.5),
        text=['# Post-hoc: Steigungen pro Gruppe',
              'emtrends(modell_int,',
              '  pairwise ~ gruppe, var = "alter")'],
        font_size=11)

    # --- Slide 18 (idx 17): Grenzen des Modells ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)

    limits = [
        ("Korrelation ≠ Kausalität", "Confounder können\nZusammenhang erklären"),
        ("Extrapolation", "Vorhersage ausserhalb\ndes Bereichs unzuverlässig"),
        ("Overfitting", "Zu viele Prädiktoren\nbei wenig Daten"),
        ("Modellvergleich", "Adj. R-Quadrat oder\nanova(modell1, modell2)"),
    ]
    for j, (title, desc) in enumerate(limits):
        y = Inches(1.6 + j * 1.35)
        add_hexagon(slide,
            left=Inches(0.5), top=y,
            width=Inches(4.0), height=Inches(1.1),
            text=[title], font_size=13, bold=True,
            fill_scheme='accent2' if j == 0 else 'accent4')
        add_text_box(slide,
            left=Inches(5.0), top=y,
            width=Inches(6.0), height=Inches(1.1),
            text=desc.split('\n'), font_size=13)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(6.0),
        width=Inches(9.0), height=Inches(0.8),
        text=["Modellwahl ist eine begründete Entscheidung!"],
        font_size=14)

    # --- Slide 19 (idx 18): Zusammenfassung Block B [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 20 (idx 19): Übungsaufgabe ---
    slide = prs.slides[19]
    clear_body_placeholder(slide)

    tasks = [
        "1. Einfache Regression: lm(wm_score ~ alter). b1 interpretieren",
        "2. Multiple Regression: lm(wm_score ~ alter + gruppe)",
        "3. Interaktion: lm(wm_score ~ alter * gruppe)",
        "4. Diagnostik: plot(modell) -- Annahmen erfüllt?",
        "5. Ergebnisabsatz mit Koeffizienten und R-Quadrat",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.6 + i * 1.1),
            width=Inches(9.5), height=Inches(0.9),
            text=[task], font_size=12)

    # --- Slide 21 (idx 20): Musterlösung ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.3), top=Inches(1.8),
        width=Inches(5.5), height=Inches(3.5),
        text=['# Einfach',
              'mod1 <- lm(wm_score ~ age, data = daten)',
              '',
              '# Multipel',
              'mod2 <- lm(wm_score ~ age + gruppe,',
              '           data = daten)',
              '',
              '# Interaktion',
              'mod3 <- lm(wm_score ~ age * gruppe,',
              '           data = daten)'],
        font_size=9)

    add_hexagon(slide,
        left=Inches(6.2), top=Inches(1.8),
        width=Inches(5.0), height=Inches(4.5),
        text=["Interpretation:", "",
              "Alter war ein signifikanter",
              "Prädiktor (b = 0.82, p = .009).",
              "Pro Lebensjahr stieg der",
              "WM-Score um 0.82 Punkte.",
              "R-Quadrat = .15.",
              "",
              "Interaktion Alter × Gruppe",
              "war nicht signifikant (p = .21)."],
        font_size=10, light=True)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
