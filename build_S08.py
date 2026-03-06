"""
Build S08: Datenvisualisierung I — Grundlagen ggplot2
Populates the scaffold PPTX with content from S08_Folieninhalte.md.
~6 R-generated plots + text/code shapes.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *

BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_08_Visualisierung1.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_08_Visualisierung1_filled.pptx"
PLOTS = BASE / "local" / "Vorlesung" / "plots"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Vier Datensätze haben denselben',
              'Mittelwert, SD und Korrelation.',
              'Sind sie gleich?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Visualisierung mit ggplot2")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie..."],
        ["Die Logik von ggplot2", "(Schichten, Mapping,", "Geome) erklären"],
        ["Geeignete Diagramm-", "typen für verschiedene", "Variablentypen wählen"],
        ["Erkennen, wann eine", "Grafik irreführend ist"],
    ]
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12, light=(i == 0))

    # --- Slide 8 (idx 7): Anscombe's Quartett — R-Plot ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S08_F08_anscombe.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 9 (idx 8): Grammar of Graphics: Schichtenprinzip ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)

    layers = [
        ("1. Daten", "Welcher Datensatz?"),
        ("2. Mapping (aes)", "Variable → visuelle Eigenschaft"),
        ("3. Geometrie", "Punkte, Balken, Boxplot..."),
        ("4. Skalen", "Achsenbereiche, Farben"),
        ("5. Koordinaten", "Kartesisch, polar..."),
        ("6. Theme", "Hintergrund, Schrift..."),
    ]
    for j, (layer, desc) in enumerate(layers):
        # Stacked, slightly offset for layer effect
        offset = j * 4
        add_hexagon(slide,
            left=Inches(0.5 + j * 0.15), top=Inches(1.6 + j * 0.85),
            width=Inches(4.0), height=Inches(0.7),
            text=[layer], font_size=12, bold=True,
            light=(j > 2))
        add_text_box(slide,
            left=Inches(5.0), top=Inches(1.6 + j * 0.85),
            width=Inches(6.0), height=Inches(0.7),
            text=[desc], font_size=13)

    # --- Slide 10 (idx 9): ggplot() + aes() ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=Inches(9.0), height=Inches(1.5),
        text=['ggplot(data = daten, aes(x = gruppe, y = wm_score))'],
        font_size=14)

    # Explanation hexagons
    add_hexagon(slide,
        left=Inches(0.5), top=Inches(3.8),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["data:", "Welcher Datensatz?"],
        font_size=13)

    add_hexagon(slide,
        left=Inches(3.5), top=Inches(3.8),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["aes():", "Welche Variable →", "welche Eigenschaft?"],
        font_size=12)

    add_hexagon(slide,
        left=Inches(6.5), top=Inches(3.8),
        width=HEX_CONTENT_W * 2, height=HEX_CONTENT_H,
        text=["Häufige Aesthetics:", "x, y, color, fill, size, shape"],
        font_size=12, light=True)

    add_rounded_rect(slide,
        left=Inches(2.0), top=Inches(5.8),
        width=Inches(8.0), height=Inches(1.0),
        text=["aes() allein erzeugt keine Grafik -- es fehlt das Geom!"],
        font_size=14)

    # --- Slide 11 (idx 10): Geome — R-Plot ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S08_F11_geome_uebersicht.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 12 (idx 11): Code lesen ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(6.5), height=Inches(2.8),
        text=['ggplot(daten,',
              '       aes(x = gruppe, y = wm_score,',
              '           fill = gruppe)) +',
              '  geom_boxplot() +',
              '  labs(x = "Gruppe", y = "WM-Score",',
              '       title = "WM nach Gruppe") +',
              '  theme_minimal()'],
        font_size=12)

    # Annotations for each line
    annotations = [
        ("Zeile 1-3:", "Daten + Mapping"),
        ("Zeile 4:", "Geometrie (Boxplot)"),
        ("Zeile 5-6:", "Beschriftungen"),
        ("Zeile 7:", "Erscheinungsbild"),
    ]
    for j, (line, desc) in enumerate(annotations):
        add_hexagon(slide,
            left=Inches(7.5), top=Inches(1.8 + j * 1.1),
            width=Inches(3.5), height=Inches(0.9),
            text=[f"{line} {desc}"], font_size=12, light=True)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(5.5),
        width=Inches(9.0), height=Inches(1.0),
        text=["Jede Zeile hat eine klare Aufgabe.",
              "Lesen und verstehen -- nicht Syntax memorieren!"],
        font_size=13)

    # --- Slide 13 (idx 12): Zusammenfassung Grammar of Graphics [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 14 (idx 13): Gut vs. schlecht — R-Plot ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S08_F14_gut_vs_schlecht.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 15 (idx 14): Entscheidungsbaum Diagrammwahl ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)

    choices = [
        ("1 Variable metrisch", "Histogramm"),
        ("1 Variable kategorial", "Balkendiagramm"),
        ("2 Var. metrisch × metrisch", "Scatterplot"),
        ("2 Var. kategorial × metrisch", "Boxplot"),
        ("2 Var. kategorial × kategorial", "Gruppierter Balken"),
    ]
    for j, (condition, chart) in enumerate(choices):
        y = Inches(1.6 + j * 1.05)
        add_hexagon(slide,
            left=Inches(0.5), top=y,
            width=Inches(5.0), height=Inches(0.85),
            text=[condition], font_size=12)
        add_arrow(slide,
            left=Inches(5.8), top=y + Emu(HEX_CONTENT_H // 4),
            width=Inches(0.8), height=ARROW_H,
            direction='right')
        add_hexagon(slide,
            left=Inches(7.0), top=y,
            width=Inches(4.0), height=Inches(0.85),
            text=[chart], font_size=13, bold=True, light=True)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(6.0),
        width=Inches(9.0), height=Inches(0.8),
        text=["Skalenniveau und Anzahl der Variablen bestimmen den Diagrammtyp"],
        font_size=13)

    # --- Slide 16 (idx 15): Eine Variable — R-Plot ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S08_F16_eine_variable.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 17 (idx 16): Zwei Variablen — R-Plot ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S08_F17_zwei_variablen.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 18 (idx 17): Irreführende Grafiken — R-Plot ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S08_F18_irrefuehrend.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 19 (idx 18): Checkliste gute Grafik [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 20 (idx 19): Übungsaufgabe ---
    slide = prs.slides[19]
    clear_body_placeholder(slide)

    tasks = [
        "1. Histogramm für wm_score: Verteilungsform?",
        "2. Boxplot wm_score nach gruppe: Was fällt auf?",
        "3. Code lesen: Was ändert fill = gruppe?",
        "Bonus: Scatterplot für zwei metrische Variablen",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.8 + i * 1.3),
            width=Inches(9.5), height=Inches(1.0),
            text=[task], font_size=13,
            light=(i == len(tasks) - 1))

    # --- Slide 21 (idx 20): Musterlösung ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.3), top=Inches(1.8),
        width=Inches(5.5), height=Inches(4.5),
        text=['# Aufgabe 1: Histogramm',
              'ggplot(daten, aes(x = wm_score)) +',
              '  geom_histogram(bins = 20)',
              '',
              '# Aufgabe 2: Boxplot',
              'ggplot(daten, aes(x = gruppe,',
              '                  y = wm_score,',
              '                  fill = gruppe)) +',
              '  geom_boxplot() +',
              '  labs(x = "Gruppe",',
              '       y = "WM-Score (Punkte)")'],
        font_size=9)

    add_hexagon(slide,
        left=Inches(6.2), top=Inches(1.8),
        width=Inches(5.0), height=Inches(3.0),
        text=["Interpretation:", "",
              "Histogramm: leicht linksschiefe",
              "Verteilung.",
              "Boxplot: EG höhere Werte als KG,",
              "aber mit grösserer Streuung."],
        font_size=11, light=True)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
