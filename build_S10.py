"""
Build S10: Ergebnisberichte in R — Quarto & Markdown
Populates the scaffold PPTX with content from S10_Folieninhalte.md.
No R-generated plots — text, code boxes, and diagram shapes only.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *


BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_10_Quarto.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_10_Quarto_filled.pptx"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.5), top=Inches(2.2),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Sie haben Ihre Analyse fertig.',
              'Ihr Ergebnis steht in der R-Konsole.',
              'Wie kommt es in den Bericht?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Quarto & Markdown")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie..."],
        ["Unterschied reproduzierbar", "vs. nicht-reproduzierbar", "dokumentiert erklären"],
        ["Ein Quarto-Dokument mit", "eingebettetem R-Code", "erstellen"],
        ["Entscheiden, was in einem", "Analysebericht nachvollziehbar", "dokumentiert sein muss"],
    ]
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12, light=(i == 0))

    # --- Slide 8 (idx 7): Das Problem: Copy-Paste-Workflow ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)

    # Flow: R → Screenshot → Word → Daten geändert → zurück zu R
    steps = [
        ("R: Ergebnis", "accent4"),
        ("Screenshot/\nKopieren", "accent4"),
        ("Word:\nEinfügen", "accent4"),
        ("Daten\ngeändert!", "accent2"),
        ("Alles\nnochmal!", "accent2"),
    ]
    for j, (label, scheme) in enumerate(steps):
        add_hexagon(slide,
            left=Inches(0.5 + j * 2.2), top=Inches(1.8),
            width=HEX_CONTENT_W, height=HEX_CONTENT_H,
            text=label.split('\n'), font_size=12, fill_scheme=scheme)
        if j < len(steps) - 1:
            add_arrow(slide,
                left=Inches(2.0 + j * 2.2), top=Inches(2.3),
                width=Inches(0.7), height=ARROW_H,
                direction='right')

    # Error sources
    problems = [
        "Falsche Version eingefügt",
        "Tabelle nicht aktualisiert",
        "Grafik aus anderer Analyse",
    ]
    for j, prob in enumerate(problems):
        add_hexagon(slide,
            left=Inches(0.8 + j * 3.5), top=Inches(4.2),
            width=Inches(3.0), height=Inches(0.9),
            text=[prob], font_size=11, fill_scheme='accent1')

    add_rounded_rect(slide,
        left=Inches(2.0), top=Inches(5.5),
        width=Inches(8.0), height=Inches(1.0),
        text=["Manuelles Übertragen ist fehleranfällig",
              "und nicht reproduzierbar"],
        font_size=14)

    # --- Slide 9 (idx 8): Die Lösung: Alles in einem Dokument ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)

    # Left: concept
    add_hexagon(slide,
        left=Inches(0.8), top=Inches(1.8),
        width=HEX_LARGE_W, height=HEX_LARGE_H,
        text=["Ein Dokument:", "Text + R-Code", "+ Output"],
        font_size=14)

    # Flow: .qmd → [Rendern] → Bericht
    add_hexagon(slide,
        left=Inches(3.5), top=Inches(4.0),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=[".qmd Datei"], font_size=13)
    add_arrow(slide,
        left=Inches(5.5), top=Inches(4.5),
        width=Inches(1.5), height=ARROW_H,
        direction='right')
    add_hexagon(slide,
        left=Inches(7.5), top=Inches(4.0),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Fertiger Bericht", "(HTML/PDF)"],
        font_size=12, fill_scheme='accent6')

    # Right: key statement
    add_rounded_rect(slide,
        left=Inches(4.0), top=Inches(1.8),
        width=Inches(7.0), height=Inches(1.5),
        text=["Der Bericht IST die Analyse",
              "-- nicht eine Kopie davon"],
        font_size=16)

    # --- Slide 10 (idx 9): Quarto: Drei Bausteine ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)

    bausteine = [
        ("YAML-Header", "Metadaten: Titel, Autor,\nDatum, Format"),
        ("Markdown", "Text: Überschriften,\nListen, Fett/Kursiv"),
        ("Code-Chunks", "R-Code: wird beim Rendern\nausgeführt"),
    ]
    for j, (title, desc) in enumerate(bausteine):
        x = Inches(0.8 + j * 3.8)
        add_hexagon(slide,
            left=x, top=Inches(2.0),
            width=HEX_LARGE_W, height=HEX_LARGE_H,
            text=[title], font_size=16, bold=True)
        add_text_box(slide,
            left=x - Inches(0.3), top=Inches(4.2),
            width=Inches(3.2), height=Inches(1.5),
            text=desc.split('\n'), font_size=13)

    # --- Slide 11 (idx 10): YAML-Header ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=Inches(4.5), height=Inches(3.0),
        text=['---',
              'title: "WM-Analyse"',
              'author: "Ihr Name"',
              'date: today',
              'format: html',
              '---'],
        font_size=14)

    # Annotations
    annotations = [
        ("title", "Titel im Bericht"),
        ("format: html", "HTML / pdf / docx"),
        ("date: today", "Automatisches Datum"),
    ]
    for j, (key, desc) in enumerate(annotations):
        add_hexagon(slide,
            left=Inches(6.5), top=Inches(1.8 + j * 1.5),
            width=HEX_CONTENT_W, height=HEX_CONTENT_H,
            text=[key, desc], font_size=12, light=True)

    add_rounded_rect(slide,
        left=Inches(2.0), top=Inches(5.5),
        width=Inches(8.0), height=Inches(1.0),
        text=["Minimal 4 Zeilen -- der Header steuert das Gesamtbild"],
        font_size=14)

    # --- Slide 12 (idx 11): Markdown-Grundlagen ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)

    # Left: Syntax
    add_text_box(slide,
        left=Inches(0.5), top=Inches(1.6),
        width=Inches(3.0), height=Inches(0.5),
        text=["Markdown-Syntax:"], font_size=14, bold=True)
    add_code_box(slide,
        left=Inches(0.5), top=Inches(2.2),
        width=Inches(4.5), height=Inches(4.0),
        text=['# Überschrift 1',
              '## Überschrift 2',
              '',
              '**fett**',
              '*kursiv*',
              '',
              '- Punkt 1',
              '- Punkt 2',
              '',
              '1. Erstens',
              '2. Zweitens',
              '',
              '[Link](url)'],
        font_size=11)

    # Right: Rendered result description
    add_text_box(slide,
        left=Inches(5.5), top=Inches(1.6),
        width=Inches(3.0), height=Inches(0.5),
        text=["Ergebnis:"], font_size=14, bold=True)
    results = [
        "Grosse Überschrift",
        "Unterüberschrift",
        "fett / kursiv",
        "Aufzählung",
        "Nummerierte Liste",
        "Hyperlink",
    ]
    for j, res in enumerate(results):
        add_hexagon(slide,
            left=Inches(5.5), top=Inches(2.2 + j * 0.7),
            width=Inches(4.0), height=Inches(0.6),
            text=[res], font_size=11, light=True)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(6.0),
        width=Inches(9.0), height=Inches(0.8),
        text=["Markdown ist einfacher als Word -- und erzeugt sauberere Dokumente"],
        font_size=13)

    # --- Slide 13 (idx 12): Quarto vs. Copy-Paste [Zwischentitelfolie] ---
    # Section slide — title already set. No additional shapes.

    # --- Slide 14 (idx 13): Leitfrage Block B ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(2.5), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Wie kommen R-Ergebnisse in den Text',
              '-- ohne Copy-Paste?'],
        font_size=16)

    add_text_box(slide,
        left=Inches(3.0), top=Inches(4.5),
        width=Inches(6.0), height=Inches(1.0),
        text=['Der Mittelwert betrug ???.'],
        font_size=18, bold=True)

    # --- Slide 15 (idx 14): Code-Chunks ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(6.5), height=Inches(2.5),
        text=['```{r}',
              '#| echo: true',
              '#| message: false',
              '',
              'daten <- read.csv2(here("daten", "WM_Fragebögen.csv"))',
              'mean(daten$wm_score, na.rm = TRUE)',
              '```'],
        font_size=12)

    # Chunk options
    options = [
        ("echo: true", "Code zeigen (Default)"),
        ("echo: false", "Code verbergen, nur Output"),
        ("message: false", "Meldungen unterdrücken"),
        ("warning: false", "Warnungen unterdrücken"),
    ]
    for j, (opt, desc) in enumerate(options):
        add_hexagon(slide,
            left=Inches(7.5), top=Inches(1.8 + j * 1.2),
            width=Inches(3.5), height=Inches(1.0),
            text=[f"#| {opt}", desc], font_size=11, light=True)

    # --- Slide 16 (idx 15): Sichtbarkeit: Wer liest den Bericht? ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)

    audiences = [
        ("Kolleg*innen", "echo: true", "Code zeigen für\nReproduzierbarkeit"),
        ("Auftraggeber", "echo: false", "Nur Ergebnisse\nzeigen"),
        ("Prüfung/Abgabe", "echo: true", "Code dokumentiert\ndie Arbeit"),
    ]
    for j, (audience, setting, reason) in enumerate(audiences):
        y = Inches(1.8 + j * 1.6)
        add_hexagon(slide,
            left=Inches(0.8), top=y,
            width=HEX_CONTENT_W, height=HEX_CONTENT_H,
            text=[audience], font_size=14, bold=True)
        add_arrow(slide,
            left=Inches(2.8), top=y + Emu(HEX_CONTENT_H // 3),
            width=Inches(1.0), height=ARROW_H,
            direction='right')
        add_code_box(slide,
            left=Inches(4.2), top=y + Emu(HEX_CONTENT_H // 6),
            width=Inches(2.5), height=Inches(0.8),
            text=[setting], font_size=12)
        add_text_box(slide,
            left=Inches(7.2), top=y,
            width=Inches(4.0), height=Inches(1.2),
            text=reason.split('\n'), font_size=13)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(5.8),
        width=Inches(9.0), height=Inches(1.0),
        text=["Die Zielgruppe bestimmt die Darstellung -- nicht den Inhalt"],
        font_size=14)

    # --- Slide 17 (idx 16): Inline R-Code ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)

    # Source code
    add_text_box(slide,
        left=Inches(0.5), top=Inches(1.6),
        width=Inches(4.0), height=Inches(0.5),
        text=["Quellcode:"], font_size=14, bold=True)
    add_code_box(slide,
        left=Inches(0.5), top=Inches(2.2),
        width=Inches(10.0), height=Inches(1.2),
        text=['Der Mittelwert betrug',
              '`r round(mean(daten$wm_score, na.rm = TRUE), 1)`',
              'Punkte.'],
        font_size=12)

    # Rendered text
    add_text_box(slide,
        left=Inches(0.5), top=Inches(3.8),
        width=Inches(4.0), height=Inches(0.5),
        text=["Gerenderter Text:"], font_size=14, bold=True)
    add_hexagon(slide,
        left=Inches(0.5), top=Inches(4.5),
        width=Inches(10.0), height=Inches(1.0),
        text=["Der Mittelwert betrug 62.3 Punkte."],
        font_size=16, fill_scheme='accent6')

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(6.0),
        width=Inches(9.0), height=Inches(0.8),
        text=["Inline-Code = der stärkste Beweis für Reproduzierbarkeit"],
        font_size=14)

    # --- Slide 18 (idx 17): Ausführungsreihenfolge ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)

    # Chunk sequence: 4 boxes stacked
    chunks = [
        ("Chunk 1: library()", "accent6"),
        ("Chunk 2: Daten importieren", "accent6"),
        ("Chunk 3: Berechnung", "accent6"),
        ("Chunk 4: Grafik", "accent6"),
    ]
    for j, (label, scheme) in enumerate(chunks):
        add_hexagon(slide,
            left=Inches(1.5), top=Inches(1.8 + j * 1.2),
            width=Inches(4.0), height=Inches(0.9),
            text=[label], font_size=13, fill_scheme=scheme)
        if j < len(chunks) - 1:
            add_arrow(slide,
                left=Inches(3.0), top=Inches(2.8 + j * 1.2),
                width=Inches(0.4), height=Inches(0.4),
                direction='down')

    # Error example
    add_hexagon(slide,
        left=Inches(6.5), top=Inches(2.0),
        width=Inches(4.5), height=Inches(1.5),
        text=["Häufiger Fehler:", "Variable in Chunk 3 erstellt,",
              "in Chunk 1 benutzt"],
        font_size=12, fill_scheme='accent2')

    add_code_box(slide,
        left=Inches(6.5), top=Inches(4.0),
        width=Inches(4.5), height=Inches(1.0),
        text=["Error: object 'daten' not found"],
        font_size=12)

    add_rounded_rect(slide,
        left=Inches(6.5), top=Inches(5.5),
        width=Inches(4.5), height=Inches(1.2),
        text=["Regel: Import & Pakete", "immer ganz oben!"],
        font_size=14)

    # --- Slide 19 (idx 18): Dokumentationsstandard ---
    slide = prs.slides[18]
    clear_body_placeholder(slide)

    checklist = [
        "1. Forschungsfrage: Was wird untersucht?",
        "2. Datenquelle & Import: Woher? Wie eingelesen?",
        "3. Aufbereitungsschritte: Was verändert -- warum?",
        "4. Analyseentscheidungen: Welche Methode? Warum?",
        "5. Ergebnisse mit Interpretation",
    ]
    for j, item in enumerate(checklist):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.6 + j * 1.1),
            width=Inches(9.5), height=Inches(0.9),
            text=[item], font_size=14)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(6.0),
        width=Inches(9.0), height=Inches(0.8),
        text=["Analyseentscheidungen sind wissenschaftliche Entscheidungen"],
        font_size=13)

    # --- Slide 20 (idx 19): Zusammenfassung Block B [Zwischentitelfolie] ---
    # Section slide — title already set.

    # --- Slide 21 (idx 20): Übungsaufgabe ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    tasks = [
        "1. Neues Quarto-Dokument (.qmd) mit YAML-Header erstellen",
        "2. Chunk 1: WM-Daten importieren (read.csv2 + here)",
        "3. Chunk 2: Deskriptive Statistik (group_by + summarise)",
        "4. Chunk 3: ggplot erstellen (Boxplot oder Histogramm)",
        "5. Inline-Code: Satz mit Mittelwert als Inline-Code",
        "Bonus: Chunk-Optionen setzen (echo, message)",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.6 + i * 0.95),
            width=Inches(9.5), height=Inches(0.8),
            text=[task], font_size=12,
            light=(i == len(tasks) - 1))

    # --- Slide 22 (idx 21): Musterlösung ---
    slide = prs.slides[21]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.3), top=Inches(1.8),
        width=Inches(5.5), height=Inches(5.0),
        text=['---',
              'title: "WM-Analyse"',
              'author: "Ihr Name"',
              'date: today',
              'format: html',
              '---',
              '',
              '## Daten importieren',
              '',
              '```{r}',
              '#| message: false',
              'library(tidyverse)',
              'library(here)',
              'daten <- read.csv2(',
              '  here("daten", "WM_Fragebögen.csv")',
              ')',
              '```'],
        font_size=9)

    add_code_box(slide,
        left=Inches(6.0), top=Inches(1.8),
        width=Inches(5.5), height=Inches(5.0),
        text=['## Deskriptive Statistik',
              '',
              '```{r}',
              'daten |>',
              '  group_by(gender) |>',
              '  summarise(',
              '    M = mean(wm_score, na.rm = TRUE),',
              '    SD = sd(wm_score, na.rm = TRUE)',
              '  )',
              '```',
              '',
              'Der Mittelwert betrug',
              '`r round(mean(daten$wm_score,',
              '  na.rm = TRUE), 1)` Punkte.'],
        font_size=9)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
