"""
Build S09: Datenvisualisierung II — Erweiterungen
Populates the scaffold PPTX with content from S09_Folieninhalte.md.
~9 R-generated plots + text/code shapes.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *

BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_09_Visualisierung2.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_09_Visualisierung2_filled.pptx"
PLOTS = BASE / "local" / "Vorlesung" / "plots"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Sie haben eine Grafik erstellt.',
              'Ihr Betreuer sagt: Nicht geeignet',
              'für ein Paper. Was fehlt?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Visualisierung II")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie..."],
        ["Mehrdimensionale Daten", "durch Facettierung", "strukturiert darstellen"],
        ["Konfidenzintervalle", "visuell korrekt", "interpretieren"],
        ["Entscheiden, wann visuelle", "Komplexität hilfreich", "oder irreführend ist"],
    ]
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12, light=(i == 0))

    # --- Slide 8 (idx 7): Progression — R-Plot ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S09_F08_progression.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 9 (idx 8): facet_wrap — R-Plot ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(1.5),
        text=['ggplot(daten, aes(x = wm_score)) +',
              '  geom_histogram() +',
              '  facet_wrap(~session)'],
        font_size=12)

    insert_image(slide, str(PLOTS / "S09_F09_facet_wrap.png"),
                 Inches(5.8), Inches(1.8), Inches(5.8), Inches(4.0))

    add_rounded_rect(slide,
        left=Inches(0.5), top=Inches(5.8),
        width=Inches(5.0), height=Inches(1.0),
        text=["Viele Gruppen, Verteilungen", "vergleichen → facet_wrap()"],
        font_size=13)

    # --- Slide 10 (idx 9): wrap vs grid — R-Plot ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S09_F10_wrap_vs_grid.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 11 (idx 10): Fehlerbalken — R-Plot ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S09_F11_fehlerbalken.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 12 (idx 11): CI-Overlap — R-Plot ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S09_F12_ci_overlap.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 13 (idx 12): Entscheidungstabelle [Zwischentitelfolie or content] ---
    slide = prs.slides[12]
    # Check if it's a section slide by counting shapes
    if len(slide.shapes) <= 2:
        # Section slide — just the title
        pass
    else:
        clear_body_placeholder(slide)

    # --- Slide 14 (idx 13): Default ugly plot — R-Plot ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S09_F14_default_ugly.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 15 (idx 14): Themes — R-Plot ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S09_F15_themes.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 16 (idx 15): Labels — R-Plot ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(2.5),
        text=['+ labs(',
              '    title = "WM-Score nach Gruppe",',
              '    x = "Versuchsgruppe",',
              '    y = "WM-Score (Punkte)",',
              '    caption = "Daten: WM-Studie 2024"',
              '  )'],
        font_size=11)

    insert_image(slide, str(PLOTS / "S09_F16_labels.png"),
                 Inches(5.8), Inches(1.8), Inches(5.8), Inches(4.0))

    add_rounded_rect(slide,
        left=Inches(0.5), top=Inches(5.5),
        width=Inches(5.0), height=Inches(1.3),
        text=["Keine Variablennamen als Labels!",
              "Einheiten in Klammern angeben."],
        font_size=13)

    # --- Slide 17 (idx 16): Patchwork — R-Plot ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(4.5), height=Inches(2.5),
        text=['library(patchwork)',
              '',
              'p1 <- ggplot(...) + geom_histogram()',
              'p2 <- ggplot(...) + geom_boxplot()',
              '',
              'p1 + p2  # nebeneinander',
              'p1 / p2  # übereinander'],
        font_size=11)

    insert_image(slide, str(PLOTS / "S09_F17_patchwork.png"),
                 Inches(5.5), Inches(1.8), Inches(6.0), Inches(4.2))

    # --- Slide 18 (idx 17): ggsave() ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=Inches(9.0), height=Inches(1.5),
        text=['ggsave("abbildungen/wm_boxplot.pdf",',
              '       width = 15, height = 10, units = "cm", dpi = 300)'],
        font_size=13)

    # Format comparison
    formats = [
        ("PDF", "Vektorgrafik, skalierbar", "Paper, Poster"),
        ("PNG", "Rastergrafik, feste Auflösung", "Präsentation, Web"),
    ]
    for j, (fmt, desc, use) in enumerate(formats):
        y = Inches(3.8 + j * 1.6)
        add_hexagon(slide,
            left=Inches(0.8), top=y,
            width=Inches(2.5), height=Inches(1.2),
            text=[fmt], font_size=16, bold=True)
        add_text_box(slide,
            left=Inches(3.8), top=y,
            width=Inches(3.5), height=Inches(1.2),
            text=[desc], font_size=13)
        add_hexagon(slide,
            left=Inches(7.5), top=y,
            width=Inches(3.5), height=Inches(1.2),
            text=[use], font_size=13, light=True)

    # --- Slide 19 (idx 18): Checkliste [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 20 (idx 19): Übungsaufgabe ---
    slide = prs.slides[19]
    clear_body_placeholder(slide)

    tasks = [
        "1. Boxplot wm_score ~ gruppe, facettiert nach session",
        "2. Fehlerbalken hinzufügen (M ± 1 SD oder 95% CI)",
        "3. Theme und Beschriftung anpassen: theme_bw(), deutsche Labels",
        "Bonus: patchwork -- zweite Grafik kombinieren",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.8 + i * 1.3),
            width=Inches(9.5), height=Inches(1.0),
            text=[task], font_size=13,
            light=(i == len(tasks) - 1))

    # --- Slide 21 (idx 20): Musterlösung ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.3), top=Inches(1.8),
        width=Inches(5.5), height=Inches(5.0),
        text=['# Grundplot + Facetten',
              'p1 <- ggplot(daten,',
              '  aes(x = gruppe, y = wm_score,',
              '      fill = gruppe)) +',
              '  geom_boxplot() +',
              '  facet_wrap(~session) +',
              '  theme_bw() +',
              '  labs(x = "Versuchsgruppe",',
              '       y = "WM-Score (Punkte)",',
              '       fill = "Gruppe")',
              '',
              '# Patchwork',
              'p2 <- ggplot(daten,',
              '  aes(x = wm_score)) +',
              '  geom_histogram(bins = 20)',
              '',
              'p1 + p2'],
        font_size=9)

    add_hexagon(slide,
        left=Inches(6.2), top=Inches(1.8),
        width=Inches(5.0), height=Inches(3.0),
        text=["Schrittweiser Aufbau:", "",
              "1. Grundplot",
              "2. + Facetten",
              "3. + Fehlerbalken",
              "4. + Theme/Labs",
              "5. + Patchwork"],
        font_size=11, light=True)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
