"""
Build S06: Grundlagen — Programmieren mit R
Populates the scaffold PPTX with content from S06_Folieninhalte.md.
All text-based — no R-generated images needed.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *

BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_06_Programmieren.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_06_Programmieren_filled.pptx"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.5), top=Inches(2.2),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Sie müssen dieselbe Datenprüfung',
              'für 5 verschiedene Datensätze durchführen.',
              'Was machen Sie?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Programmieren mit R")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie…"],
        ["Begründen, warum eigene", "Funktionen Reproduzierbar-", "keit erhöhen"],
        ["Eine einfache Analyse-", "funktion schreiben (Input →", "Verarbeitung → Output)"],
        ["Entscheiden, wann eine", "Funktion sinnvoll ist vs.", "direkte Ausführung"],
    ]
    # Honeycomb layout: 1 intro hex top-center, 3 objective hexes below
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12,
                    fill_scheme='accent4' if i > 0 else 'accent4',
                    light=(i == 0))

    # --- Slide 8 (idx 7): Warum Funktionen? Das Copy-Paste-Problem ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)

    # Left: Problem
    add_hexagon(slide,
        left=Inches(0.8), top=Inches(1.8),
        width=HEX_LARGE_W, height=HEX_LARGE_H,
        text=["Problem:", "Gleicher Code 5× kopiert", "→ 1 Fehler = 5 Korrekturen"],
        font_size=14, bold=False,
        fill_scheme='accent2')

    # Right: Solution
    add_hexagon(slide,
        left=Inches(6.5), top=Inches(1.8),
        width=HEX_LARGE_W, height=HEX_LARGE_H,
        text=["Lösung:", "DRY — Don't Repeat Yourself", "1 Funktion → 5 Aufrufe"],
        font_size=14, bold=False)

    # Arrow between
    add_arrow(slide,
        left=Inches(3.5), top=Inches(2.5),
        width=Inches(2.5), height=ARROW_H,
        direction='right')

    # Bottom takeaway
    add_rounded_rect(slide,
        left=Inches(2.0), top=Inches(4.5),
        width=Inches(8.0), height=Inches(1.2),
        text=["Funktionen = dokumentierte, testbare,",
              "wiederverwendbare Bausteine"],
        font_size=16)

    # --- Slide 9 (idx 8): Anatomie einer R-Funktion ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)

    # Code example
    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(2.5),
        text=['pruefe_bereich <- function(x, min = 0, max = 100) {',
              '  ausserhalb <- sum(x < min | x > max, na.rm = TRUE)',
              '  return(ausserhalb)',
              '}'],
        font_size=12)

    # Flow diagram: Input → Verarbeitung → Output
    flow_y = Inches(5.0)
    add_hexagon(slide,
        left=Inches(0.8), top=flow_y,
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Input:", "x, min, max"],
        font_size=12)
    add_arrow(slide,
        left=Inches(3.0), top=Inches(5.5),
        width=Inches(1.0), height=ARROW_H,
        direction='right')
    add_hexagon(slide,
        left=Inches(4.5), top=flow_y,
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Verarbeitung:", "Prüfung & Zählung"],
        font_size=12)
    add_arrow(slide,
        left=Inches(6.7), top=Inches(5.5),
        width=Inches(1.0), height=ARROW_H,
        direction='right')
    add_hexagon(slide,
        left=Inches(8.2), top=flow_y,
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Output:", "Anzahl ausserhalb"],
        font_size=12)

    # --- Slide 10 (idx 9): Argumente: Pflicht vs. Standard ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)

    # Annotated function call
    add_code_box(slide,
        left=Inches(1.5), top=Inches(1.8),
        width=Inches(9.0), height=Inches(1.5),
        text=['pruefe_bereich(x = daten$wm_score, min = 10, max = 90)'],
        font_size=14)

    # Annotation hexagons
    add_hexagon(slide,
        left=Inches(1.0), top=Inches(3.8),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Pflichtargument:", "x hat keinen Default"],
        font_size=12, fill_scheme='accent1')

    add_hexagon(slide,
        left=Inches(4.2), top=Inches(3.8),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Standardwert:", "min = 0, max = 100", "(überschrieben)"],
        font_size=12, fill_scheme='accent3')

    add_hexagon(slide,
        left=Inches(7.5), top=Inches(3.8),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Kernaussage:", "Standardwerte", "dokumentieren Annahmen"],
        font_size=12)

    # --- Slide 11 (idx 10): Beispiel: WM-Daten validieren ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)

    # Function definition (left)
    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(3.0),
        text=['pruefe_bereich <- function(x, min = 0, max = 100) {',
              '  n_aussen <- sum(x < min | x > max, na.rm = TRUE)',
              '  n_na <- sum(is.na(x))',
              '  cat("Ausserhalb:", n_aussen,',
              '      "/ Fehlend:", n_na, "\\n")',
              '}'],
        font_size=11)

    # Call + output (right)
    add_code_box(slide,
        left=Inches(6.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(1.5),
        text=['pruefe_bereich(daten$wm_score)',
              '# Ausserhalb: 2 / Fehlend: 5'],
        font_size=12)

    add_rounded_rect(slide,
        left=Inches(6.5), top=Inches(4.0),
        width=Inches(5.0), height=Inches(1.5),
        text=["Funktion dokumentiert Kriterium (0–100),",
              "ist wiederverwendbar, und gibt",
              "strukturierten Output"],
        font_size=12)

    # --- Slide 12 (idx 11): Entscheidung: Funktion oder direkt? ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)

    # Decision tree as hexagon cascade
    decisions = [
        ("Code >2× verwendet?", "→ Funktion schreiben"),
        ("Komplexe Logik?", "→ Funktion zur Dokumentation"),
        ("Flexibilität nötig?", "→ Funktion mit Argumenten"),
        ("Einmalig & einfach?", "→ Direkt ausführen"),
    ]
    for i, (question, answer) in enumerate(decisions):
        y = Inches(1.8) + Inches(i * 1.4)
        add_hexagon(slide,
            left=Inches(1.5), top=y,
            width=HEX_CONTENT_W, height=HEX_CONTENT_H,
            text=[question], font_size=13)
        add_text_box(slide,
            left=Inches(4.5), top=y + Emu(HEX_CONTENT_H // 4),
            width=Inches(5.0), height=Inches(0.6),
            text=[answer], font_size=16, bold=True)

    # --- Slide 13 (idx 12): Funktionen: Drei Kernpunkte [Zwischentitelfolie] ---
    # Already has correct title and layout. This is a section slide.
    # No additional content needed — title is sufficient.

    # --- Slide 14 (idx 13): Automatisierung: Warum nicht einzeln? ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['10 Variablen auf fehlende Werte prüfen.',
              'Einzeln tippen — oder automatisieren?'],
        font_size=16)

    # Table header simulation
    cols = ["wm_score", "alter", "stress", "motivation", "schlaf", "..."]
    for j, col in enumerate(cols):
        add_hexagon(slide,
            left=Inches(0.5 + j * 1.9), top=Inches(4.5),
            width=Inches(1.7), height=Inches(0.9),
            text=[col], font_size=11, light=True)

    # --- Slide 15 (idx 14): for-Schleifen ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(2.5),
        text=['for (var in names(daten)) {',
              '  n_na <- sum(is.na(daten[[var]]))',
              '  cat(var, ":", n_na, "fehlend\\n")',
              '}'],
        font_size=12)

    # Output
    add_code_box(slide,
        left=Inches(6.5), top=Inches(1.8),
        width=Inches(4.5), height=Inches(2.0),
        text=['wm_score : 5 fehlend',
              'alter    : 0 fehlend',
              'stress   : 12 fehlend',
              '...'],
        font_size=12)

    # Loop diagram
    add_hexagon(slide,
        left=Inches(2.0), top=Inches(5.0),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Sequenz:", "var1, var2, var3, ..."],
        font_size=12, light=True)
    add_arrow(slide,
        left=Inches(4.2), top=Inches(5.5),
        width=Inches(1.5), height=ARROW_H,
        direction='right')
    add_hexagon(slide,
        left=Inches(6.0), top=Inches(5.0),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["für jedes Element:", "Code ausführen"],
        font_size=12)

    # --- Slide 16 (idx 15): if-else ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(6.0), height=Inches(2.8),
        text=['anteil_na <- mean(is.na(daten$wm_score))',
              '',
              'if (anteil_na > 0.10) {',
              '  cat("Warnung: >10% fehlend!\\n")',
              '} else {',
              '  cat("OK: Fehlquote akzeptabel\\n")',
              '}'],
        font_size=12)

    # Flowchart: diamond question → two paths
    add_hexagon(slide,
        left=Inches(7.5), top=Inches(1.8),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["anteil_na", "> 0.10 ?"],
        font_size=13, fill_scheme='accent3')
    add_hexagon(slide,
        left=Inches(6.5), top=Inches(4.0),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Ja →", "Warnung ausgeben"],
        font_size=12, fill_scheme='accent1')
    add_hexagon(slide,
        left=Inches(8.8), top=Inches(4.0),
        width=HEX_CONTENT_W, height=HEX_CONTENT_H,
        text=["Nein →", "OK melden"],
        font_size=12, fill_scheme='accent6')

    # --- Slide 17 (idx 16): Schleife + Bedingung kombinieren ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=Inches(7.5), height=Inches(3.8),
        text=['for (var in names(daten)) {',
              '  if (is.numeric(daten[[var]])) {',
              '    anteil <- mean(is.na(daten[[var]]))',
              '    if (anteil > 0.10) {',
              '      cat("⚠", var, ":",',
              '          round(anteil * 100), "% fehlend\\n")',
              '    }',
              '  }',
              '}'],
        font_size=12)

    # Color legend
    labels = [
        ("Schleife", "accent3"),
        ("Typprüfung", "accent4"),
        ("NA-Check", "accent1"),
        ("Warnung", "accent2"),
    ]
    for j, (label, scheme) in enumerate(labels):
        add_hexagon(slide,
            left=Inches(9.0), top=Inches(1.8 + j * 1.2),
            width=Inches(1.8), height=Inches(0.9),
            text=[label], font_size=11, fill_scheme=scheme)

    # --- Slide 18 (idx 17): apply & map ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)

    # Left: for-loop
    add_text_box(slide,
        left=Inches(0.5), top=Inches(1.7),
        width=Inches(3.0), height=Inches(0.5),
        text=["for-Schleife (4 Zeilen):"], font_size=14, bold=True)
    add_code_box(slide,
        left=Inches(0.5), top=Inches(2.3),
        width=Inches(5.0), height=Inches(2.2),
        text=['for (var in names(daten)) {',
              '  cat(var,',
              '      sum(is.na(daten[[var]])),',
              '      "\\n")',
              '}'],
        font_size=11)

    # Right: sapply / map
    add_text_box(slide,
        left=Inches(6.0), top=Inches(1.7),
        width=Inches(5.0), height=Inches(0.5),
        text=["Funktionale Alternative:"], font_size=14, bold=True)
    add_code_box(slide,
        left=Inches(6.0), top=Inches(2.3),
        width=Inches(5.5), height=Inches(1.0),
        text=['sapply(daten, function(x) sum(is.na(x)))'],
        font_size=12)
    add_code_box(slide,
        left=Inches(6.0), top=Inches(3.5),
        width=Inches(5.5), height=Inches(1.0),
        text=['map(daten, ~sum(is.na(.x)))'],
        font_size=12)

    # Bottom takeaway
    add_rounded_rect(slide,
        left=Inches(2.0), top=Inches(5.2),
        width=Inches(8.0), height=Inches(1.0),
        text=["Gleiche Logik, unterschiedliche Schreibweise.",
              "map() ist der tidyverse-Weg."],
        font_size=14)

    # --- Slide 19 (idx 18): Entscheidungshilfe [Zwischentitelfolie] ---
    # Section slide — title already set. No additional shapes needed.

    # --- Slide 20 (idx 19): Übung ---
    slide = prs.slides[19]
    clear_body_placeholder(slide)

    tasks = [
        "1. Funktion pruefe_variable(): Wertebereich prüfen + Anzahl fehlender Werte ausgeben",
        "2. Funktion mit for-Schleife oder map() auf alle numerischen Spalten anwenden",
        "3. Erweitern: Warnung wenn >10% fehlend",
        "Bonus: Funktion für Anteil ausserhalb eines Bereichs",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.8 + i * 1.3),
            width=Inches(9.0), height=Inches(1.0),
            text=[task], font_size=13,
            light=(i == len(tasks) - 1))

    # --- Slide 21 (idx 20): Musterlösung ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(4.8),
        text=['# Aufgabe 1: Funktion definieren',
              'pruefe_variable <- function(x, min = 0, max = 100) {',
              '  n_aussen <- sum(x < min | x > max, na.rm = TRUE)',
              '  n_na <- sum(is.na(x))',
              '  anteil_na <- mean(is.na(x))',
              '  cat("Ausserhalb:", n_aussen,',
              '      "/ Fehlend:", n_na, "\\n")',
              '  if (anteil_na > 0.10) {',
              '    cat("⚠ Warnung: >10% fehlend\\n")',
              '  }',
              '}'],
        font_size=10)

    add_code_box(slide,
        left=Inches(6.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(3.0),
        text=['# Aufgabe 2: Auf alle Spalten anwenden',
              '# Mit for-Schleife:',
              'for (var in names(daten)) {',
              '  if (is.numeric(daten[[var]])) {',
              '    cat(var, ": ")',
              '    pruefe_variable(daten[[var]])',
              '  }',
              '}',
              '',
              '# Mit map:',
              'map(daten[num_vars], pruefe_variable)'],
        font_size=10)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
