"""
Reusable shape library for Statistik IV slide production.
Matches the exact visual language of sessions 1-5:
  - Hexagons: accent4 navy fill, white text, Arial 14pt, centered
  - Arrows: pink fill, no outline
  - Rounded rects: grey fill, bold text
  - Code boxes: Consolas font, accent4 fill
"""

from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn, nsmap

# Namespace map for building XML elements
NSMAP_A = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
NSMAP_P = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
NSMAP_R = {'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}


# ---------- XML building helpers ----------

def _make_element(tag, attrib=None, text=None):
    """Create an lxml element with the given qualified tag name."""
    el = etree.SubElement(etree.Element('dummy'), qn(tag))
    el = deepcopy(el)
    if attrib:
        for k, v in attrib.items():
            el.set(k, str(v))
    if text is not None:
        el.text = text
    return el


def _build_style_xml():
    """Build the standard shape style XML (matches existing slides exactly)."""
    style = _make_element('p:style')

    lnRef = etree.SubElement(style, qn('a:lnRef'), idx='2')
    clr1 = etree.SubElement(lnRef, qn('a:schemeClr'), val='accent1')
    etree.SubElement(clr1, qn('a:shade'), val='50000')

    fillRef = etree.SubElement(style, qn('a:fillRef'), idx='1')
    etree.SubElement(fillRef, qn('a:schemeClr'), val='accent1')

    effectRef = etree.SubElement(style, qn('a:effectRef'), idx='0')
    etree.SubElement(effectRef, qn('a:schemeClr'), val='accent1')

    fontRef = etree.SubElement(style, qn('a:fontRef'), idx='minor')
    etree.SubElement(fontRef, qn('a:schemeClr'), val='lt1')

    return style


def _build_txBody(lines, font_size=1400, font_name=None, bold=False, lang='de-DE'):
    """Build txBody XML with centered text, 117% line spacing.

    Args:
        lines: list of strings (one per paragraph), or a single string.
        font_size: in hundredths of a point (1400 = 14pt).
        font_name: explicit font name (e.g., 'Consolas'). None = theme font (Arial).
        bold: whether text is bold.
    """
    if isinstance(lines, str):
        lines = [lines]

    txBody = _make_element('p:txBody')
    bodyPr = etree.SubElement(txBody, qn('a:bodyPr'),
                              lIns='72000', tIns='72000', rIns='72000', bIns='72000',
                              rtlCol='0', anchor='ctr', anchorCtr='0')
    etree.SubElement(txBody, qn('a:lstStyle'))

    for line in lines:
        p = etree.SubElement(txBody, qn('a:p'))
        pPr = etree.SubElement(p, qn('a:pPr'), algn='ctr')
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        etree.SubElement(lnSpc, qn('a:spcPct'), val='117000')

        r = etree.SubElement(p, qn('a:r'))
        rPr_attrib = {'lang': lang, 'sz': str(font_size), 'dirty': '0'}
        if bold:
            rPr_attrib['b'] = '1'
        rPr = etree.SubElement(r, qn('a:rPr'), **rPr_attrib)
        if font_name:
            etree.SubElement(rPr, qn('a:latin'), typeface=font_name)
            etree.SubElement(rPr, qn('a:cs'), typeface=font_name)

        t = etree.SubElement(r, qn('a:t'))
        t.text = line

    return txBody


def _build_spPr(left, top, width, height, prst, fill_scheme='accent4',
                fill_mods=None, no_line=True):
    """Build spPr XML for a preset geometry shape.

    Args:
        fill_scheme: scheme color name (e.g., 'accent4', 'accent2').
        fill_mods: list of (mod_name, value) tuples, e.g., [('lumMod', '60000'), ('lumOff', '40000')].
    """
    spPr = _make_element('p:spPr')

    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    etree.SubElement(xfrm, qn('a:off'), x=str(left), y=str(top))
    etree.SubElement(xfrm, qn('a:ext'), cx=str(width), cy=str(height))

    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'), prst=prst)
    etree.SubElement(prstGeom, qn('a:avLst'))

    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    clr = etree.SubElement(solidFill, qn('a:schemeClr'), val=fill_scheme)
    if fill_mods:
        for mod_name, mod_val in fill_mods:
            etree.SubElement(clr, qn(f'a:{mod_name}'), val=str(mod_val))

    ln = etree.SubElement(spPr, qn('a:ln'))
    if no_line:
        etree.SubElement(ln, qn('a:noFill'))

    return spPr


def _add_shape(slide, nvSpPr_name, spPr, style, txBody):
    """Add a complete sp (shape) element to a slide's spTree."""
    spTree = slide._element.find(qn('p:cSld') + '/' + qn('p:spTree'))

    sp = etree.SubElement(spTree, qn('p:sp'))

    # nvSpPr
    nvSpPr = etree.SubElement(sp, qn('p:nvSpPr'))
    # Generate a unique ID
    existing_ids = [int(s.get('id', 0))
                    for s in spTree.findall('.//' + qn('p:cNvPr'))]
    new_id = max(existing_ids, default=1) + 1
    cNvPr = etree.SubElement(nvSpPr, qn('p:cNvPr'), id=str(new_id), name=nvSpPr_name)
    cNvSpPr = etree.SubElement(nvSpPr, qn('p:cNvSpPr'))
    nvPr = etree.SubElement(nvSpPr, qn('p:nvPr'))

    sp.append(spPr)
    sp.append(style)
    if txBody is not None:
        sp.append(txBody)

    return sp, new_id


# ---------- Public API ----------

def add_hexagon(slide, left, top, width, height, text,
                font_size=14, light=False, fill_scheme='accent4', bold=False):
    """Add a hexagon shape matching the Statistik IV visual language.

    Args:
        slide: pptx Slide object.
        left, top, width, height: position/size in EMU. Use Inches() or Emu() helpers.
        text: string or list of strings (one per line inside the hexagon).
        font_size: in points (default 14).
        light: if True, uses lighter fill variant (lumMod=60000, lumOff=40000).
        fill_scheme: override fill color scheme (default 'accent4' = navy).
        bold: if True, text is bold.

    Returns:
        (shape_element, shape_id) tuple.
    """
    fill_mods = [('lumMod', '60000'), ('lumOff', '40000')] if light else None

    spPr = _build_spPr(left, top, width, height, 'hexagon',
                       fill_scheme=fill_scheme, fill_mods=fill_mods)
    style = _build_style_xml()
    txBody = _build_txBody(text, font_size=font_size * 100, bold=bold)

    return _add_shape(slide, 'Sechseck', spPr, style, txBody)


def add_arrow(slide, left, top, width, height, direction='right',
              fill_scheme='accent1', fill_mods=None):
    """Add an arrow shape.

    Args:
        direction: 'right', 'down', 'left', or 'up'.
        fill_scheme: scheme color (default 'accent1' = magenta/pink).
    """
    prst_map = {
        'right': 'rightArrow',
        'down': 'downArrow',
        'left': 'leftArrow',
        'up': 'upArrow',
    }
    prst = prst_map.get(direction, 'rightArrow')

    spPr = _build_spPr(left, top, width, height, prst,
                       fill_scheme=fill_scheme, fill_mods=fill_mods)
    style = _build_style_xml()
    txBody = _build_txBody('', font_size=1400)

    return _add_shape(slide, f'Pfeil nach {direction}', spPr, style, txBody)


def add_rounded_rect(slide, left, top, width, height, text,
                     font_size=14, bold=True):
    """Add a rounded rectangle with grey fill, matching existing style.

    Uses bg1 at lumMod=50000 (medium grey).
    """
    fill_mods = [('lumMod', '50000')]
    spPr = _build_spPr(left, top, width, height, 'roundRect',
                       fill_scheme='bg1', fill_mods=fill_mods)
    style = _build_style_xml()
    txBody = _build_txBody(text, font_size=font_size * 100, bold=bold)

    return _add_shape(slide, 'Abgerundetes Rechteck', spPr, style, txBody)


def add_code_box(slide, left, top, width, height, text, font_size=12):
    """Add a text box with Consolas font for R code display.

    Uses accent4 fill (matching hexagon style) with Consolas monospace.
    """
    spPr = _build_spPr(left, top, width, height, 'rect', fill_scheme='accent4')
    style = _build_style_xml()

    # For code: left-aligned, not centered
    txBody = _build_txBody(text, font_size=font_size * 100, font_name='Consolas')
    # Override alignment to left
    for pPr in txBody.findall('.//' + qn('a:pPr')):
        pPr.set('algn', 'l')
    # Override body anchor to top
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 't')

    return _add_shape(slide, 'Code', spPr, style, txBody)


def add_text_box(slide, left, top, width, height, text,
                 font_size=14, bold=False, color_scheme=None):
    """Add a plain text box (no fill, no border).

    Args:
        color_scheme: text color scheme (None = default black via 'tx1').
    """
    spPr = _make_element('p:spPr')
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    etree.SubElement(xfrm, qn('a:off'), x=str(left), y=str(top))
    etree.SubElement(xfrm, qn('a:ext'), cx=str(width), cy=str(height))
    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'), prst='rect')
    etree.SubElement(prstGeom, qn('a:avLst'))
    etree.SubElement(spPr, qn('a:noFill'))
    ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))

    style = _build_style_xml()
    txBody = _build_txBody(text, font_size=font_size * 100, bold=bold)

    # Override text color if needed
    if color_scheme:
        for rPr in txBody.findall('.//' + qn('a:rPr')):
            solidFill = etree.SubElement(rPr, qn('a:solidFill'))
            etree.SubElement(solidFill, qn('a:schemeClr'), val=color_scheme)
    else:
        # Default: black text for text boxes (override the lt1 from style)
        for rPr in txBody.findall('.//' + qn('a:rPr')):
            solidFill = etree.SubElement(rPr, qn('a:solidFill'))
            etree.SubElement(solidFill, qn('a:schemeClr'), val='tx1')

    return _add_shape(slide, 'Textfeld', spPr, style, txBody)


def insert_image(slide, image_path, left, top, width, height):
    """Insert a PNG image at the specified position.

    Returns the picture shape.
    """
    return slide.shapes.add_picture(str(image_path), left, top, width, height)


def clear_body_placeholder(slide, placeholder_idx=13):
    """Clear the body text placeholder (idx=13 for 'Titel und Text, 1-spaltig')."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == placeholder_idx:
            shape.text_frame.clear()
            return True
    return False


def set_title(slide, text):
    """Set the title placeholder (idx=0) text."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            shape.text_frame.clear()
            shape.text_frame.paragraphs[0].text = text
            return True
    return False


# ---------- Standard dimensions (EMU) ----------

# Content area boundaries (below title bar)
CONTENT_LEFT = Inches(0.5)
CONTENT_TOP = Inches(1.6)
CONTENT_WIDTH = Inches(11.0)
CONTENT_HEIGHT = Inches(5.4)

# Hexagon sizes
HEX_LERNZIELE_W = 1729632   # ~1.81"
HEX_LERNZIELE_H = 1554480   # ~1.63"
HEX_CONTENT_W = 1627889     # ~1.71"
HEX_CONTENT_H = 1463040     # ~1.53"
HEX_LARGE_W = 1936750       # ~2.03"
HEX_LARGE_H = 1740624       # ~1.82"

# Arrow sizes
ARROW_W = Inches(1.0)
ARROW_H = Inches(0.4)

# Image area (full width below title)
IMG_LEFT = Inches(0.8)
IMG_TOP = Inches(1.8)
IMG_W = Inches(10.4)
IMG_H = Inches(5.2)

# Slide dimensions
SLIDE_W = Inches(12.0)
SLIDE_H = Inches(7.5)
