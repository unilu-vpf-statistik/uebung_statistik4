"""
Build S11: Anwendung I — Gruppenvergleiche & ANOVA
Populates the scaffold PPTX with content from S11_Folieninhalte.md.
3 R-generated plots + text/code shapes.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *

BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_11_ANOVA.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_11_ANOVA_filled.pptx"
PLOTS = BASE / "local" / "Vorlesung" / "plots"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['KG hat M = 58, EG hat M = 65.',
              'Unterscheiden sich die Gruppen',
              'wirklich? Was brauchen Sie?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Gruppenvergleiche & ANOVA")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie..."],
        ["ANOVA als Varianz-", "zerlegungsmodell (nicht", "als Black-Box) erklären"],
        ["Modellannahmen prüfen", "und Konsequenzen von", "Verletzungen beurteilen"],
        ["Signifikanz von Effekt-", "grösse unterscheiden", "und begründet berichten"],
    ]
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12, light=(i == 0))

    # --- Slide 8 (idx 7): Gruppenunterschied — R-Plot ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S11_F08_boxplots_gruppen.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 9 (idx 8): Varianzzerlegung ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=["Gesamtvarianz (SS_total) =",
              "Zwischen (SS_between) +",
              "Innerhalb (SS_within)"],
        font_size=14)

    # Two scenarios
    add_hexagon(slide,
        left=Inches(0.8), top=Inches(4.0),
        width=Inches(4.5), height=Inches(2.5),
        text=["Szenario 1:", "Viel SS_between,", "wenig SS_within",
              "→ Klarer Gruppenunterschied"],
        font_size=12, fill_scheme='accent6')

    add_hexagon(slide,
        left=Inches(6.0), top=Inches(4.0),
        width=Inches(4.5), height=Inches(2.5),
        text=["Szenario 2:", "Wenig SS_between,", "viel SS_within",
              "→ Unterschied im Rauschen"],
        font_size=12, fill_scheme='accent2')

    add_rounded_rect(slide,
        left=Inches(5.5), top=Inches(1.8),
        width=Inches(5.5), height=Inches(1.5),
        text=["Frage: Wie viel Varianz wird", "durch die Gruppen erklärt?"],
        font_size=14)

    # --- Slide 10 (idx 9): F-Statistik — R-Plot ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S11_F10_f_verteilung.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 11 (idx 10): Effektgrössen ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)

    metrics = [
        ("Eta-Quadrat (eta sq.)", "SS_between / SS_total",
         "Anteil erklärter Varianz", "eta sq. = .12 → 12%"),
        ("Cohen's d", "(M1 - M2) / SD_pooled",
         "Standardisierter Unterschied", "d = 0.2/0.5/0.8"),
    ]
    for j, (name, formula, concept, example) in enumerate(metrics):
        y = Inches(1.8 + j * 2.5)
        add_hexagon(slide,
            left=Inches(0.5), top=y,
            width=Inches(3.0), height=Inches(2.0),
            text=[name], font_size=14, bold=True)
        add_text_box(slide,
            left=Inches(4.0), top=y,
            width=Inches(3.5), height=Inches(2.0),
            text=[formula, concept], font_size=13)
        add_hexagon(slide,
            left=Inches(8.0), top=y,
            width=Inches(3.5), height=Inches(2.0),
            text=[example], font_size=13, light=True)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(5.8),
        width=Inches(9.0), height=Inches(1.0),
        text=["Immer Effektgrösse berichten -- nicht nur p!"],
        font_size=14)

    # --- Slide 12 (idx 11): R-Output lesen ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(3.0),
        text=['library(afex)',
              '',
              'anova_ergebnis <- aov_ez(',
              '  id = "id",',
              '  dv = "wm_score",',
              '  between = "gruppe",',
              '  data = daten',
              ')'],
        font_size=11)

    # Output columns explanation
    columns = [
        ("Effect", "Welcher Faktor?"),
        ("df", "Freiheitsgrade"),
        ("F", "F-Statistik"),
        ("ges", "Eta-Quadrat (Effektgrösse)"),
        ("p", "p-Wert"),
    ]
    for j, (col, desc) in enumerate(columns):
        add_hexagon(slide,
            left=Inches(6.0), top=Inches(1.8 + j * 1.0),
            width=Inches(5.0), height=Inches(0.8),
            text=[f"{col}: {desc}"], font_size=12, light=True)

    add_rounded_rect(slide,
        left=Inches(6.0), top=Inches(5.8),
        width=Inches(5.0), height=Inches(1.0),
        text=["Lesen und interpretieren,", "nicht auswendig lernen!"],
        font_size=13)

    # --- Slide 13 (idx 12): Zusammenfassung Block A [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 14 (idx 13): Leitfrage Block B ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.0),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Darf ich die ANOVA rechnen?',
              'Was muss ich vorher prüfen?'],
        font_size=16)

    checks = ["Normalverteilung", "Varianzhomogenität", "Unabhängigkeit"]
    for j, check in enumerate(checks):
        add_hexagon(slide,
            left=Inches(1.5 + j * 3.5), top=Inches(4.5),
            width=Inches(3.0), height=Inches(1.2),
            text=[check], font_size=14, light=True)

    # --- Slide 15 (idx 14): Modellannahmen — R-Plot ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S11_F15_qq_plots.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 16 (idx 15): Konsequenzen von Verletzungen ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)

    violations = [
        ("Normalverteilung", "Robust bei ähnlichen N.\nSonst: Kruskal-Wallis", "accent3"),
        ("Varianzhomogenität", "Welch-ANOVA als\nAlternative (afex Default)", "accent3"),
        ("Unabhängigkeit", "Gravierendste Verletzung!\n→ Anderes Modell nötig", "accent2"),
    ]
    for j, (name, consequence, scheme) in enumerate(violations):
        y = Inches(1.8 + j * 1.6)
        add_hexagon(slide,
            left=Inches(0.5), top=y,
            width=HEX_CONTENT_W, height=HEX_CONTENT_H,
            text=[name], font_size=14, bold=True)
        add_arrow(slide,
            left=Inches(2.5), top=y + Emu(HEX_CONTENT_H // 3),
            width=Inches(1.0), height=ARROW_H,
            direction='right')
        add_hexagon(slide,
            left=Inches(4.0), top=y,
            width=Inches(7.0), height=HEX_CONTENT_H,
            text=consequence.split('\n'), font_size=12, fill_scheme=scheme)

    # --- Slide 17 (idx 16): Post-hoc ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.0), height=HEX_LARGE_H,
        text=["ANOVA sagt nur:", "Irgendein Unterschied existiert.",
              "Welche Gruppen? → Post-hoc"],
        font_size=13)

    add_hexagon(slide,
        left=Inches(6.0), top=Inches(1.8),
        width=Inches(5.0), height=HEX_LARGE_H,
        text=["Warum nicht viele t-Tests?",
              "10 Tests: P(mind. 1 Fehler) ≈ 40%!",
              "→ Alpha-Korrektur nötig"],
        font_size=13, fill_scheme='accent2')

    add_code_box(slide,
        left=Inches(1.0), top=Inches(4.5),
        width=Inches(9.0), height=Inches(1.5),
        text=['library(emmeans)',
              'emmeans(anova_ergebnis, pairwise ~ gruppe,',
              '        adjust = "bonferroni")'],
        font_size=12)

    # --- Slide 18 (idx 17): R-Workflow kompakt ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)

    workflow = [
        ("1. Daten vorbereiten", "Tidy format, Faktor kodiert"),
        ("2. Annahmen prüfen", "QQ-Plot, Levene-Test"),
        ("3. ANOVA rechnen", "aov_ez(id, dv, between, data)"),
        ("4. Effektgrösse", "ges (eta sq.) aus Output"),
        ("5. Post-hoc", "emmeans(..., pairwise ~ gruppe)"),
    ]
    for j, (step, detail) in enumerate(workflow):
        add_hexagon(slide,
            left=Inches(0.5), top=Inches(1.6 + j * 1.1),
            width=Inches(3.5), height=Inches(0.9),
            text=[step], font_size=13, bold=True)
        add_code_box(slide,
            left=Inches(4.5), top=Inches(1.6 + j * 1.1),
            width=Inches(6.5), height=Inches(0.9),
            text=[detail], font_size=11)

    add_rounded_rect(slide,
        left=Inches(1.5), top=Inches(6.2),
        width=Inches(9.0), height=Inches(0.8),
        text=["Erst Annahmen, dann Analyse, dann Effektgrösse!"],
        font_size=14)

    # --- Slide 19 (idx 18): Zusammenfassung Block B [Zwischentitelfolie] ---
    # Section slide.

    # --- Slide 20 (idx 19): Übungsaufgabe ---
    slide = prs.slides[19]
    clear_body_placeholder(slide)

    tasks = [
        "1. Annahmen prüfen: QQ-Plot + Levene-Test",
        "2. ANOVA: aov_ez() mit wm_score als AV, gruppe als UV",
        "3. Output lesen: F, eta sq., Signifikanz?",
        "4. Post-hoc: emmeans() -- welche Gruppen?",
        "5. Ergebnissatz mit Effektgrösse formulieren",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.6 + i * 1.1),
            width=Inches(9.5), height=Inches(0.9),
            text=[task], font_size=12)

    # --- Slide 21 (idx 20): Musterlösung ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.3), top=Inches(1.8),
        width=Inches(5.5), height=Inches(3.5),
        text=['# 1. Annahmen',
              'shapiro.test(daten$wm_score[daten$gruppe == "KG"])',
              'car::leveneTest(wm_score ~ gruppe, data = daten)',
              '',
              '# 2. ANOVA',
              'anova_ergebnis <- aov_ez(',
              '  id = "id", dv = "wm_score",',
              '  between = "gruppe", data = daten)'],
        font_size=9)

    add_hexagon(slide,
        left=Inches(6.2), top=Inches(1.8),
        width=Inches(5.0), height=Inches(4.5),
        text=["Ergebnis:", "",
              "Die ANOVA zeigte einen",
              "signifikanten Gruppenunterschied",
              "F(1, 98) = 8.42, p = .004,",
              "eta sq. = .08.",
              "",
              "Die EG erzielte signifikant",
              "höhere Werte als die KG",
              "(Diff = 6.9, 95% CI [2.1, 11.7])"],
        font_size=10, light=True)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
