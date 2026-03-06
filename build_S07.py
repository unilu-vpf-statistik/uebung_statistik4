"""
Build S07: Deskriptive Statistik berechnen
Populates the scaffold PPTX with content from S07_Folieninhalte.md.
4 R-generated plots + text/code shapes.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx_helpers import *

BASE = Path(__file__).parent
SCAFFOLD = BASE / "local" / "Vorlesung" / "FS26_Statistik4_07_DeskriptiveStatistik.pptx"
OUTPUT = BASE / "local" / "Vorlesung" / "FS26_Statistik4_07_DeskriptiveStatistik_filled.pptx"
PLOTS = BASE / "local" / "Vorlesung" / "plots"


def build():
    prs = Presentation(str(SCAFFOLD))

    # --- Slide 3 (idx 2): Aktivierung ---
    slide = prs.slides[2]
    clear_body_placeholder(slide)
    set_title(slide, "Aktivierung")
    add_hexagon(slide,
        left=Inches(2.0), top=Inches(2.2),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=['Welcher Kennwert beschreibt die',
              'typische Leistung im WM-Test',
              'am besten -- und warum?'],
        font_size=16)

    # --- Slide 6 (idx 5): Lernziele ---
    slide = prs.slides[5]
    clear_body_placeholder(slide)
    set_title(slide, "Lernziele: Deskriptive Statistik")

    lernziele = [
        ["Nach dieser Vorlesung", "sollten Sie..."],
        ["Begründen, welche Kenn-", "werte welche Forschungs-", "frage beantworten"],
        ["Gruppierte Zusammen-", "fassungen in R umsetzen", "und interpretieren"],
        ["Robuste vs. nicht-robuste", "Masse unterscheiden und", "kontextabhängig einsetzen"],
    ]
    x_positions = [Inches(4.5), Inches(1.5), Inches(4.5), Inches(7.5)]
    y_positions = [Inches(1.8), Inches(3.8), Inches(3.8), Inches(3.8)]
    for i, (lz, x, y) in enumerate(zip(lernziele, x_positions, y_positions)):
        add_hexagon(slide, left=x, top=y,
                    width=HEX_LERNZIELE_W, height=HEX_LERNZIELE_H,
                    text=lz, font_size=12, light=(i == 0))

    # --- Slide 8 (idx 7): Rolle der Deskriptivstatistik ---
    slide = prs.slides[7]
    clear_body_placeholder(slide)

    add_hexagon(slide,
        left=Inches(1.0), top=Inches(1.8),
        width=HEX_LARGE_W * 2, height=HEX_LARGE_H,
        text=["Deskriptivstatistik ist kein", "Selbstzweck -- sie ist", "Entscheidungsgrundlage"],
        font_size=14)

    # Process diagram: steps
    steps = ["Import", "Validierung", "Aufbereitung", "Deskription", "Modellierung"]
    for j, step in enumerate(steps):
        scheme = 'accent1' if step == "Deskription" else 'accent4'
        fs = 13 if step == "Deskription" else 12
        add_hexagon(slide,
            left=Inches(0.5 + j * 2.2), top=Inches(4.5),
            width=HEX_CONTENT_W, height=HEX_CONTENT_H,
            text=[step], font_size=fs, fill_scheme=scheme,
            bold=(step == "Deskription"))
        if j < len(steps) - 1:
            add_arrow(slide,
                left=Inches(2.0 + j * 2.2), top=Inches(5.0),
                width=Inches(0.7), height=ARROW_H,
                direction='right')

    # --- Slide 9 (idx 8): Lagemaße — R-Plot ---
    slide = prs.slides[8]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S07_F09_histogramme.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 10 (idx 9): Streuungsmaße — R-Plot ---
    slide = prs.slides[9]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S07_F10_streuung.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 11 (idx 10): Skalenniveau → Kennwert ---
    slide = prs.slides[10]
    clear_body_placeholder(slide)

    rows = [
        ("Nominal", "Modus", "Häufigkeiten", "gruppe"),
        ("Ordinal", "Median", "IQR", "Schulnoten"),
        ("Metrisch\n(symmetrisch)", "Mittelwert", "SD", "wm_score"),
        ("Metrisch\n(schief)", "Median", "IQR", "Reaktionszeiten"),
    ]
    # Header
    headers = ["Skalenniveau", "Lage", "Streuung", "Beispiel"]
    for j, h in enumerate(headers):
        add_hexagon(slide,
            left=Inches(0.3 + j * 2.8), top=Inches(1.7),
            width=Inches(2.5), height=Inches(0.7),
            text=[h], font_size=12, bold=True)

    for i, (niveau, lage, streuung, beispiel) in enumerate(rows):
        y = Inches(2.7 + i * 1.15)
        vals = [niveau, lage, streuung, beispiel]
        for j, val in enumerate(vals):
            add_hexagon(slide,
                left=Inches(0.3 + j * 2.8), top=y,
                width=Inches(2.5), height=Inches(0.9),
                text=val.split('\n'), font_size=11, light=True)

    # --- Slide 12 (idx 11): WM-Scores beschreiben — R-Plot ---
    slide = prs.slides[11]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S07_F12_wm_histogram.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 13 (idx 12): Zusammenfassung Kennwerte [Zwischentitelfolie] ---
    # Section slide — title already set.

    # --- Slide 14 (idx 13): Leitfrage Block B — R-Plot ---
    slide = prs.slides[13]
    clear_body_placeholder(slide)
    insert_image(slide, str(PLOTS / "S07_F14_density_gruppen.png"),
                 IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    # --- Slide 15 (idx 14): group_by + summarise ---
    slide = prs.slides[14]
    clear_body_placeholder(slide)

    # Split-Apply-Combine as hexagons
    sac = [("Split:", "Daten nach\nGruppe aufteilen"),
           ("Apply:", "Kennwert pro\nGruppe berechnen"),
           ("Combine:", "Ergebnisse\nzusammenfügen")]
    for j, (title, desc) in enumerate(sac):
        add_hexagon(slide,
            left=Inches(0.5 + j * 3.8), top=Inches(1.8),
            width=HEX_LARGE_W, height=HEX_LARGE_H,
            text=[title] + desc.split('\n'), font_size=13)
        if j < len(sac) - 1:
            add_arrow(slide,
                left=Inches(2.8 + j * 3.8), top=Inches(2.5),
                width=Inches(1.5), height=ARROW_H,
                direction='right')

    add_code_box(slide,
        left=Inches(1.0), top=Inches(4.5),
        width=Inches(6.0), height=Inches(1.8),
        text=['daten |>',
              '  group_by(gruppe) |>',
              '  summarise(M = mean(wm_score, na.rm = TRUE))'],
        font_size=12)

    add_hexagon(slide,
        left=Inches(7.5), top=Inches(4.5),
        width=Inches(3.5), height=Inches(1.8),
        text=["Ergebnis:", "KG: M = 58.2", "EG: M = 65.1"],
        font_size=13, light=True)

    # --- Slide 16 (idx 15): Mehrere Kennwerte ---
    slide = prs.slides[15]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(6.0), height=Inches(3.5),
        text=['daten |>',
              '  group_by(gruppe) |>',
              '  summarise(',
              '    M   = mean(wm_score, na.rm = TRUE),',
              '    SD  = sd(wm_score, na.rm = TRUE),',
              '    Mdn = median(wm_score, na.rm = TRUE),',
              '    N   = n()',
              '  )'],
        font_size=11)

    add_hexagon(slide,
        left=Inches(7.0), top=Inches(1.8),
        width=Inches(4.5), height=Inches(2.0),
        text=["M(EG) > M(KG), aber", "SD(EG) auch grösser", "→ mehr Streuung in EG"],
        font_size=13, light=True)

    add_rounded_rect(slide,
        left=Inches(7.0), top=Inches(4.2),
        width=Inches(4.5), height=Inches(1.2),
        text=["Mdn ≈ M → symmetrische", "Verteilungen in beiden Gruppen"],
        font_size=13)

    # --- Slide 17 (idx 16): across() ---
    slide = prs.slides[16]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(6.0), height=Inches(3.0),
        text=['daten |>',
              '  group_by(gruppe) |>',
              '  summarise(across(',
              '    where(is.numeric),',
              '    list(M = mean, SD = sd),',
              '    .names = "{.col}_{.fn}"',
              '  ))'],
        font_size=11)

    # Tidyselect helpers
    helpers = [
        ("where(is.numeric)", "Alle numerischen Spalten"),
        ('starts_with("wm")', 'Spalten die mit "wm" beginnen'),
        ("c(alter, stress)", "Bestimmte Spalten"),
    ]
    for j, (code, desc) in enumerate(helpers):
        add_hexagon(slide,
            left=Inches(7.0), top=Inches(1.8 + j * 1.4),
            width=Inches(4.5), height=Inches(1.1),
            text=[code, desc], font_size=11, light=True)

    # --- Slide 18 (idx 17): Häufigkeitstabellen ---
    slide = prs.slides[17]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(5.0), height=Inches(1.2),
        text=['# Einfache Häufigkeit',
              'count(daten, gruppe)'],
        font_size=12)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(3.3),
        width=Inches(5.0), height=Inches(1.2),
        text=['# Kreuztabelle',
              'table(daten$gruppe, daten$geschlecht)'],
        font_size=12)

    add_code_box(slide,
        left=Inches(0.5), top=Inches(4.8),
        width=Inches(5.0), height=Inches(1.5),
        text=['# Mit Prozenten',
              'janitor::tabyl(daten, gruppe,',
              '               geschlecht)'],
        font_size=12)

    # Example cross table
    add_hexagon(slide,
        left=Inches(6.0), top=Inches(2.0),
        width=Inches(5.0), height=Inches(3.5),
        text=["Kreuztabelle:", "",
              "       m    f    Total",
              "KG    52   57    109",
              "EG    35   37     72",
              "Total 87   94    181"],
        font_size=12, light=True)

    # --- Slide 19 (idx 18): Zusammenfassung Deskriptiver Workflow [Zwischentitelfolie] ---
    # Section slide — title already set.

    # --- Slide 20 (idx 19): Übungsaufgabe ---
    slide = prs.slides[19]
    clear_body_placeholder(slide)

    tasks = [
        "1. M, SD, Mdn und N für wm_score, gruppiert nach gruppe",
        "2. across() für alle numerischen Variablen",
        "3. Kreuztabelle gruppe × geschlecht mit janitor::tabyl()",
        "4. Interpretation: Was sagen die Deskriptivstatistiken?",
        "Bonus: M vs. Mdn je Gruppe -- Hinweise auf Schiefe?",
    ]
    for i, task in enumerate(tasks):
        add_hexagon(slide,
            left=Inches(1.0), top=Inches(1.6 + i * 1.1),
            width=Inches(9.5), height=Inches(0.9),
            text=[task], font_size=12,
            light=(i == len(tasks) - 1))

    # --- Slide 21 (idx 20): Musterlösung ---
    slide = prs.slides[20]
    clear_body_placeholder(slide)

    add_code_box(slide,
        left=Inches(0.3), top=Inches(1.8),
        width=Inches(5.5), height=Inches(3.5),
        text=['daten |>',
              '  group_by(gruppe) |>',
              '  summarise(',
              '    M   = mean(wm_score, na.rm = TRUE),',
              '    SD  = sd(wm_score, na.rm = TRUE),',
              '    Mdn = median(wm_score, na.rm = TRUE),',
              '    N   = n()',
              '  )'],
        font_size=10)

    add_hexagon(slide,
        left=Inches(6.2), top=Inches(1.8),
        width=Inches(5.0), height=Inches(3.5),
        text=["Interpretation:", "",
              "Die EG zeigt einen höheren",
              "mittleren WM-Score (M = 65.1)",
              "als die KG (M = 58.2).",
              "Die Streuung ist in beiden",
              "Gruppen ähnlich (SD ≈ 12-13)."],
        font_size=11, light=True)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.name}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
